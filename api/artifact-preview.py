import io
import json
import os
import urllib.error
import urllib.parse
import urllib.request
from http.server import BaseHTTPRequestHandler

from openpyxl import load_workbook
from pptx import Presentation

MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_SHEETS = 12
MAX_ROWS = 120
MAX_COLS = 40
MAX_SLIDES = 80
MAX_TEXT_BLOCKS = 80


def _clean(value, limit=500):
    return str(value or "").strip()[:limit]


def _json_response(handler, status, payload):
    raw = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    handler.send_response(status)
    handler.send_header("Content-Type", "application/json; charset=utf-8")
    handler.send_header("Access-Control-Allow-Origin", "*")
    handler.send_header("Access-Control-Allow-Headers", "authorization, content-type")
    handler.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
    handler.send_header("Cache-Control", "no-store")
    handler.send_header("Content-Length", str(len(raw)))
    handler.end_headers()
    handler.wfile.write(raw)


def _supabase_config():
    url = os.environ.get("VITE_SUPABASE_URL") or os.environ.get("SUPABASE_URL")
    anon = os.environ.get("VITE_SUPABASE_ANON_KEY") or os.environ.get("SUPABASE_ANON_KEY")
    if not url or not anon:
        raise RuntimeError("Supabase preview auth configuration is missing.")
    return url.rstrip("/"), anon


def _verify_supabase_user(authorization):
    if not authorization or not authorization.lower().startswith("bearer "):
        raise PermissionError("Authorization bearer token is required.")
    supabase_url, anon_key = _supabase_config()
    request = urllib.request.Request(
        f"{supabase_url}/auth/v1/user",
        headers={"Authorization": authorization, "apikey": anon_key},
        method="GET",
    )
    try:
        with urllib.request.urlopen(request, timeout=8) as response:
            payload = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as error:
        raise PermissionError(f"Supabase user verification failed ({error.code}).") from error
    if not _clean(payload.get("id"), 160):
        raise PermissionError("Supabase user verification returned no user id.")


def _validate_artifact_url(value):
    supabase_url, _ = _supabase_config()
    expected = urllib.parse.urlparse(supabase_url)
    parsed = urllib.parse.urlparse(_clean(value, 5000))
    if parsed.scheme != "https" or parsed.hostname != expected.hostname:
        raise ValueError("Artifact preview URL host is invalid.")
    if not parsed.path.startswith("/storage/v1/object/sign/assistant-files/"):
        raise ValueError("Artifact preview URL path is invalid.")
    return parsed.geturl()


def _download(url):
    request = urllib.request.Request(url, headers={"User-Agent": "JetWorkArtifactPreview/1.0"}, method="GET")
    with urllib.request.urlopen(request, timeout=20) as response:
        length = int(response.headers.get("Content-Length") or "0")
        if length > MAX_FILE_BYTES:
            raise ValueError("Artifact exceeds preview size limit.")
        raw = response.read(MAX_FILE_BYTES + 1)
    if len(raw) > MAX_FILE_BYTES:
        raise ValueError("Artifact exceeds preview size limit.")
    return raw


def _cell_value(value):
    if value is None:
        return ""
    if isinstance(value, (str, int, float, bool)):
        return str(value)
    return _clean(value, 2000)


def _xlsx_preview(raw):
    workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=False)
    sheet_count = len(workbook.sheetnames)
    sheets = []
    try:
        for worksheet in workbook.worksheets[:MAX_SHEETS]:
            rows = []
            for row_index, row in enumerate(
                worksheet.iter_rows(min_row=1, max_row=MAX_ROWS, max_col=MAX_COLS, values_only=True),
                start=1,
            ):
                values = [_cell_value(value) for value in row]
                while values and values[-1] == "":
                    values.pop()
                if values:
                    rows.append({"row": row_index, "values": values})
            sheets.append({
                "name": worksheet.title,
                "rows": rows,
                "maxRow": min(int(worksheet.max_row or 0), MAX_ROWS),
                "maxColumn": min(int(worksheet.max_column or 0), MAX_COLS),
            })
    finally:
        workbook.close()
    return {"kind": "spreadsheet", "sheets": sheets, "truncated": sheet_count > MAX_SHEETS}


def _shape_text(shape):
    if getattr(shape, "has_text_frame", False):
        parts = [paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
        return "\n".join(parts)
    if getattr(shape, "has_table", False):
        lines = []
        for row in shape.table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                lines.append(" | ".join(values))
        return "\n".join(lines)
    return ""


def _pptx_preview(raw):
    presentation = Presentation(io.BytesIO(raw))
    slides = []
    total_slides = len(presentation.slides)
    for index, slide in enumerate(presentation.slides, start=1):
        if index > MAX_SLIDES:
            break
        blocks = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text:
                continue
            blocks.append({
                "text": text[:8000],
                "left": int(getattr(shape, "left", 0) or 0),
                "top": int(getattr(shape, "top", 0) or 0),
                "width": int(getattr(shape, "width", 0) or 0),
                "height": int(getattr(shape, "height", 0) or 0),
            })
            if len(blocks) >= MAX_TEXT_BLOCKS:
                break
        title = ""
        if getattr(slide.shapes, "title", None) is not None:
            title = _clean(slide.shapes.title.text, 500)
        if not title and blocks:
            title = blocks[0]["text"].split("\n", 1)[0][:500]
        slides.append({"number": index, "title": title or f"Slayt {index}", "blocks": blocks})
    return {
        "kind": "presentation",
        "slides": slides,
        "slideWidth": int(presentation.slide_width or 0),
        "slideHeight": int(presentation.slide_height or 0),
        "truncated": total_slides > MAX_SLIDES,
    }


class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            content_length = int(self.headers.get("Content-Length") or "0")
            if content_length <= 0 or content_length > 64_000:
                return _json_response(self, 413, {"error": "Preview request is empty or too large."})
            authorization = self.headers.get("Authorization") or ""
            _verify_supabase_user(authorization)
            payload = json.loads(self.rfile.read(content_length).decode("utf-8"))
            if not isinstance(payload, dict):
                return _json_response(self, 400, {"error": "Preview request must be a JSON object."})
            url = _validate_artifact_url(payload.get("url"))
            name = _clean(payload.get("name"), 500).lower()
            mime_type = _clean(payload.get("mimeType"), 500).lower()
            raw = _download(url)
            if name.endswith(".xlsx") or "spreadsheetml" in mime_type:
                preview = _xlsx_preview(raw)
            elif name.endswith(".pptx") or "presentationml" in mime_type:
                preview = _pptx_preview(raw)
            else:
                return _json_response(self, 400, {"error": "Structured preview is not supported for this file type."})
            return _json_response(self, 200, {"name": name, "byteSize": len(raw), **preview})
        except PermissionError as error:
            return _json_response(self, 401, {"error": str(error)})
        except ValueError as error:
            return _json_response(self, 400, {"error": str(error)})
        except Exception as error:
            return _json_response(self, 500, {"error": f"Artifact preview failed: {str(error)[:1000]}"})

    def do_OPTIONS(self):
        self.send_response(204)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Headers", "authorization, content-type")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.end_headers()
